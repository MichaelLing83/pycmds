#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from __future__ import annotations

from collections import defaultdict
from collections.abc import Callable
import os
from pathlib import Path
from typing import Dict, Generator, Set, Tuple, Type

from loguru import logger
import magic
from pptx import Presentation
import pptx
import pptx.shapes
import pptx.shapes.autoshape
import pptx.shapes.base

CODEC_BINARY: str = "binary"


class FileTypeCodec(object):
    """A class to handle file type encoding and decoding."""

    magic: magic.Magic = magic.Magic(mime=True, mime_encoding=True)
    MAGIC_RESULT_SEPARATOR: str = ";"

    type_codec_stat: Dict[str, int] = defaultdict(int)
    type_stat: Dict[str, int] = defaultdict(int)
    codec_stat: Dict[str, int] = defaultdict(int)

    history: Dict[Path, Tuple[str, str]] = dict()

    @classmethod
    def magic_from_file(cls, fpath: Path) -> Tuple[str, str] | None:
        fpath = fpath.resolve()
        if not fpath.exists():
            return None
        if fpath in cls.history:
            return cls.history[fpath]
        _type_codec: str = cls.magic.from_file(fpath)
        cls.type_codec_stat[_type_codec] += 1
        _type, _codec = _type_codec.split(cls.MAGIC_RESULT_SEPARATOR)
        _type = _type.strip()
        _codec = _codec.replace("charset=", "").strip()
        cls.type_stat[_type] += 1
        cls.codec_stat[_codec] += 1
        cls.history[fpath] = (_type, _codec)
        return _type, _codec

    @classmethod
    def is_binary(cls, fpath: Path) -> bool | None:
        """Check if the file is binary."""
        _type_codec: Tuple[str, str] | None = cls.magic_from_file(fpath)
        if _type_codec is None:
            return None
        _, _codec = _type_codec
        if _codec == CODEC_BINARY:
            return True
        return False

    @classmethod
    def is_text(cls, fpath: Path) -> bool | None:
        """Check if the file is text."""
        _is_binary: bool | None = cls.is_binary(fpath)
        if _is_binary is None:
            return None
        return not _is_binary

    @classmethod
    def get_codec(cls, fpath: Path) -> str | None:
        """Get the codec of the file."""
        _type_codec: Tuple[str, str] | None = cls.magic_from_file(fpath)
        if _type_codec is None:
            return None
        _, _codec = _type_codec
        return _codec

    @classmethod
    def get_type(cls, fpath: Path) -> str | None:
        """Get the type of the file."""
        _type_codec: Tuple[str, str] | None = cls.magic_from_file(fpath)
        if _type_codec is None:
            return None
        _type, _ = _type_codec
        return _type

    @classmethod
    def reset_stats(cls) -> None:
        """Reset the statistics."""
        cls.type_codec_stat.clear()
        cls.type_stat.clear()
        cls.codec_stat.clear()

    @classmethod
    def reset_history(cls) -> None:
        """Reset the history."""
        cls.history.clear()



class FileReader(object):
    def __init__(self) -> None:
        raise NotImplementedError("This class should not be instantiated directly.")

    @classmethod
    def can_read(cls, fpath: Path) -> bool:
        """Check if the file can be read."""
        raise NotImplementedError("This method should be implemented by subclasses.")

    @classmethod
    def read(cls, fpath: Path) -> Generator[str, None, None]:
        """Read the file."""
        raise NotImplementedError("This method should be implemented by subclasses.")
    
    @classmethod
    def get_reader(cls, fpath: Path) -> Type[FileReader] | None:
        if not fpath.exists():
            return None
        for _cls in cls.__subclasses__():
            if _cls.can_read(fpath):
                return _cls
        return None


class TextFileReader(FileReader):
    """A class to read text files."""

    @classmethod
    def can_read(cls, fpath: Path) -> bool:
        _ok: bool | None = FileTypeCodec.is_text(fpath)
        if _ok is None:
            return False
        return _ok

    @classmethod
    def read(cls, fpath: Path) -> Generator[str, None, None]:
        """Read a text file line by line."""
        if not fpath.exists():
            raise FileNotFoundError(f"File {fpath} does not exist.")
        _encoding: str | None = FileTypeCodec.get_codec(fpath)
        if _encoding is None:
            raise ValueError(f"Cannot determine encoding for {fpath}.")
        with open(fpath, "r", encoding=_encoding) as _f:
            for _line in _f:
                yield _line

class PptxFileReader(FileReader):
    """A class to read PowerPoint files."""
    
    @classmethod
    def can_read(cls, fpath: Path) -> bool:
        _type: str | None = FileTypeCodec.get_type(fpath)
        if _type is None:
            return False
        return _type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    @classmethod
    def read(cls, fpath: Path) -> Generator[str, None, None]:
        """Read a PowerPoint file."""
        _presentation = Presentation(str(fpath))
        # text_runs will be populated with a list of strings,
        # one for each text run in presentation
        for slide in _presentation.slides:
            for shape in slide.shapes:
                if isinstance(shape, pptx.shapes.autoshape.Shape) and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            yield run.text

if __name__ == "__main__":
    import argparse

    _arg_parser = argparse.ArgumentParser(
        description="Get file type and codec from a file."
    )
    _arg_parser.add_argument(
        "root", type=Path, help="Root dir to search files to test."
    )
    args = _arg_parser.parse_args()
    try:
        if not args.root.exists():
            logger.error("{} does not exist.", args.root)
            exit(1)
        for _dir_str, _, _fnames in os.walk(args.root):
            _dir: Path = Path(_dir_str)
            for _fname in _fnames:
                _fpath: Path = _dir / _fname
                _reader: Type[FileReader] | None = FileReader.get_reader(_fpath)
                if _reader is None:
                    continue
                else:
                    logger.debug("Found reader for {}: {}", _fpath, _reader.__name__)
        logger.info("Type codec stat: {}", FileTypeCodec.type_codec_stat)
        logger.info("Type stat: {}", FileTypeCodec.type_stat)
        logger.info("Codec stat: {}", FileTypeCodec.codec_stat)
    except KeyboardInterrupt:
        logger.info("Type codec stat: {}", FileTypeCodec.type_codec_stat)
        logger.info("Type stat: {}", FileTypeCodec.type_stat)
        logger.info("Codec stat: {}", FileTypeCodec.codec_stat)
